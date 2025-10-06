import numpy as np
import open3d as o3d
from collections import defaultdict, deque 
import os
import vtk
import sys
import math
import cv2
from vtk.util import numpy_support as nps   
import trimesh
import trimesh.transformations as tf
import pyrender 
from pptx.util import Cm
import pytesseract
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from docx import Document

''' 
INSTALL THESE IN TERMINAL:
pip install numpy
pip install open3d
pip install vtk
pip install opencv-python
pip install trimesh
pip install pyrender
pip install python-pptx
pip install pytesseract
pip install pandas
pip install python-docx
'''

pytesseract.pytesseract.tesseract_cmd = "/usr/local/bin/tesseract"



case_number = input("Enter Case Number:").strip()
case_file = f"{case_number}-downloads"
base_path = "/Users/anshshetty/Desktop/Enhatch Internship/preop planning automation" #EDIT BASE_PATH HERE, THIS IS THE PATH TO WHERE ALL THE CASE FILES ARE STORED ON YOUR COMPUTER, ALL OTHER PATHS WILL BE UPDATED AS NEEDED.
case_path = os.path.join(base_path, case_file)
femur_path = os.path.join(case_path, "fem")
femur_aligned_path = os.path.join(femur_path, "aligned")
mesh_path_femur =  os.path.join(femur_aligned_path, f"{case_number}_fem_par-cut-aligned.obj")
tibia_path = os.path.join(case_path, "tib")
tibia_aligned_path = os.path.join(tibia_path, "aligned")
mesh_path_tibia =  os.path.join(tibia_aligned_path, f"{case_number}_tib_par-cut-aligned.obj")
xray_image_path = os.path.join(case_path, "xray")
ppt_path =  os.path.join(base_path,"03-SP-xxxx_Surgical-Plan-Template_rev5.pptx")

excel_path = os.path.join(case_path, f"{case_number}-run_details.xlsx")
femur_df = pd.read_excel(excel_path, sheet_name = 0, skiprows = 3, index_col = 0)
tibia_df = pd.read_excel(excel_path, sheet_name = 1, skiprows = 3, index_col = 0)
df = pd.read_excel(excel_path, sheet_name = 0, header = None)
number_bit = case_number[:7]
word_path = os.path.join(case_path, f"{number_bit}_FRM-100-04 Case Order Checklist_rev8.docx")
doc = Document(word_path)
table = doc.tables[0]
data = []
for row in table.rows:
    data.append([cell.text.strip() for cell in row.cells])
word_df = pd.DataFrame(data)

case_id = df.iloc[1,1]
femoral_raw = femur_df.loc["Femoral Size", "Recommended"]
femoral_size = str(femoral_raw).replace("SZ_", "")
if femoral_size.isdigit():
        femoral_size = f"{int(femoral_size)/10:.1f}"
else:
    femoral_size = "NA" 
tibial_raw = tibia_df.loc["Tibial Size", "Recommended"]

#### COLOUR RESECTIONS RED #### - CHECK MESH BEFORE HAND, IF THERE ARE SPIKES ETC, ANSWER YES WHEN PROMPTED IF MESH IS DAMAGED, OTHERWISE NO.

def load_and_stats(path):
    m = o3d.io.read_triangle_mesh(path)
    if m.is_empty():
        raise ValueError(f"Could not read mesh: {path}")

    # light, safe cleanup
    m.remove_duplicated_vertices()
    m.remove_degenerate_triangles()
    m.remove_duplicated_triangles()
    m.remove_non_manifold_edges()

    V = np.asarray(m.vertices)
    T = np.asarray(m.triangles)

    # overall scale (bounding-box diagonal)
    bb_min, bb_max = V.min(axis=0), V.max(axis=0)
    diag = float(np.linalg.norm(bb_max - bb_min))

    # median edge length (typical triangle resolution)
    edges = np.vstack([T[:,[0,1]], T[:,[1,2]], T[:,[2,0]]]) #picks edge with vertices v0 and v1 for every triangle.
    edges.sort(axis=1)                      # undirected
    edges = np.unique(edges, axis=0) #picks edge with vertices v1 and v2 for every triangle.
    elens = np.linalg.norm(V[edges[:,0]] - V[edges[:,1]], axis=1) #picks edge with vertices v1 and v2 for every triangle.
    med_edge = float(np.median(elens)) if len(elens) else float('nan')

    print(f"Vertices: {len(V):,}   Triangles: {len(T):,}")
    print(f"BBox diagonal: {diag:.3f} (mesh units)")
    print(f"Median edge:   {med_edge:.3f} (mesh units)")
    return m, {"diag": diag, "med_edge": med_edge, "nV": len(V), "nF": len(T)}

# --- use it ---


def build_face_adjacency(triangles: np.ndarray):
    """
for every triangle f, we want a list of neighbor faces that share an edge with f. 
    """
    edge_to_faces = defaultdict(list)

    # Map each undirected edge -> list of incident faces
    for f, (a,b,c) in enumerate(triangles):
        for u, v in ((a,b), (b,c), (c,a)):
            if u > v:  # make edge undirected (u < v)
                u, v = v, u
            edge_to_faces[(u, v)].append(f)

    # Build adjacency only where two faces share an edge
    nbrs = [[] for _ in range(len(triangles))] #create a list of triangles where they share an edge.
    for faces in edge_to_faces.values(): #loop through every edge's list of incident faces
        if len(faces) == 2: #only processes edges that belong to exactly 2 triangles
            f0, f1 = faces
            nbrs[f0].append(f1)
            nbrs[f1].append(f0)
    return nbrs

# ---- run + quick sanity

def cluster_by_normal(triangles, tri_normals, adj, angle_threshold_degree=3.0):
    """
    BFS over face graph; grow a cluster when |dot(n_f, n_g)| >= cos(thresh).
    Returns a list of lists of face indices. Each face appears in exactly one cluster.
    """
    cos_thresh = np.cos(np.deg2rad(angle_threshold_degree)) #convert degrees to radians and then take the cosine. The dot product of 2 unit normals = the cosine of the angle between them. So if dot(x,y) >= the cos_thresh, the angle (x,y) <= threshold
    F = len(triangles)
    seen = np.zeros(F, dtype=bool)
    clusters = []

    for start in range(F): #basically goes through all the faces, figures out who shares an edge with whom. Finds normals. Starts at the first unseen face and the queue takes the next face whose neighbours they haven't explored yet, adds newly accepted neighbours to the back of the line.  
        if seen[start]:
            continue
        q = deque([start])
        seen[start] = True
        cluster = [start]

        while q:
            f = q.popleft() #keep expanding the cluster, take the next face from the front of the queue whose neighbours we'll now test
            nf = tri_normals[f]
            for g in adj[f]: #look at the faces that share an edge
                if not seen[g]: #considers faces that haven't been seen yet
                    ng = tri_normals[g] #finds normals of the faces that aren't already assigned
                    if np.dot(nf, ng) >= cos_thresh: #angle test
                        seen[g] = True #this face is now seen
                        q.append(g) #adds it to the queue so that its neighbours can be seem
                        cluster.append(g) #records g as part of the current cluster

        clusters.append(cluster)   # <-- append AFTER BFS finishes this component
    return clusters


def triangle_area(v0, v1, v2):
    return 0.5 * np.linalg.norm(np.cross(v1 - v0, v2 - v0)) #finds the area of a triangle in 3D space from its 3 corner points. Reject tiny clusters and only keep the ones large enough to matter. The cross product formula of a x b is mod(a) . mod(b) . sin(theta). Which is exactly the parallelogram area formula. And the triangle area is half of that. 

def fit_plane_get_residuals(points: np.ndarray):
    """
    Best-fit plane for a set of 3D points via SVD of centered coordinates.
    Returns (unit normal n, centroid c, max_abs_dev, rms_dev).
    If the fit is degenerate, returns (None, c, inf, inf) so caller can skip.
    """
    if points.ndim != 2 or points.shape[1] != 3 or points.shape[0] < 3:
        # not enough points for a plane
        c = points.mean(axis=0) if points.size else np.zeros(3)
        return None, c, float("inf"), float("inf") 

    # remove duplicates to avoid rank-1 pathologies
    pts = np.unique(points, axis=0)
    c = pts.mean(axis=0) #centering at the mean
    X = pts - c

    # need at least 3 unique points not collinear
    if pts.shape[0] < 3:
        return None, c, float("inf"), float("inf")

    try:
        _, _, vh = np.linalg.svd(X, full_matrices=False)  #SVD of the centered data, rows of vh are the right singular vectors, which gives the points along the direction with the least variance, equivalent to the normal of the plane.  
    except np.linalg.LinAlgError: 
        return None, c, float("inf"), float("inf")

    n = vh[-1] #the vector of least variance. This is because the SVD sorts singular values in descending order. vh[0] is the direction of largest variance in the data, vh[1] is the second-largest, vh[2] i.e. vh[-1] is the direction of smallest variance.
    n_norm = np.linalg.norm(n) #normalises the values
    if n_norm == 0:
        return None, c, float("inf"), float("inf")
    n = n / n_norm

    dists = np.abs(X @ n) #finds the signed distance to the plane that fits the cluster, at each point x in the cluster
    max_dev = float(dists.max()) if dists.size else float("inf") #finds the worst deviation of a point on the cluster, suggests whether the patch actually fits. 
    rms_dev = float(np.sqrt((dists**2).mean())) if dists.size else float("inf") #gives the mean
    return n, c, max_dev, rms_dev


def select_planar_clusters(mesh, clusters, *, med_edge, max_dev_factor=0.5, area_min_factor=3.0, rms_factor = 0.2):
    """
    Keep clusters that are big enough and planar enough, using fit_plane_get_residuals.
    """
    V = np.asarray(mesh.vertices)
    T = np.asarray(mesh.triangles)

    max_dev = max_dev_factor * med_edge #planarity tolerance, distance to plane of point
    area_min = area_min_factor * (med_edge ** 2) #minimum area for a cluster to be considered

    selected_faces = []
    kept_ct = 0

    print("\n--- Cluster diagnostics ---")
    for idx, cl in enumerate(clusters):
        # area
        area = 0.0
        for f in cl: #loop through each face index f in the cluster
            i, j, k = T[f] #gives the 3 vertex indices (i, j, k) of that face
            area += triangle_area(V[i], V[j], V[k]) #finds the sum of the areas of all the triangles in the cluster. the V[i] etc. are the 3D coordinates of those points. 

        verts = np.unique(T[cl].ravel()) #collects all vertices used by the faces in this cluster
        pts = V[verts] #looks up the 3D coordinates of those vertices, now we have a point cloud representing this cluster. 
        n, c, max_dev_meas, rms_dev_meas = fit_plane_get_residuals(pts) #fits the best plane to pts 

        rms_max = rms_factor * med_edge  # try 0.15–0.25×
        keep = (area >= area_min) and (n is not None) \
        and (max_dev_meas <= max_dev) and (rms_dev_meas <= rms_max) #keeps this cluster if it's big enough, if the fit isn't degenerate and it's within the max deviation gate. 
        print(f"#{idx:04d} faces={len(cl):5d} area={area:9.4f} "
            f"max_dev={max_dev_meas:7.4f} thresh={max_dev:7.4f} keep={keep}") #tunes the threshold

        if keep:
            selected_faces.extend(cl)
            kept_ct += 1 #accumulate all faces from this good cluster, count how many clusters passed. 

    print(f"Kept {kept_ct} / {len(clusters)} clusters")
    return np.array(sorted(set(selected_faces)), dtype=int)


# ---------- color the mesh ----------
# Open3D colors are per-vertex; we'll color vertices belonging to selected faces red.

def process_color_mesh(mesh_path, aligned_dir, case_number, label, *, angle_threshold_degree=6.0, max_dev_factor=0.2, area_min_factor=50.0, rms_factor=0.20, bone_col=(0.89, 0.855, 0.79), red_col=(1.0, 0.0, 0.0)):
    
    m, stats = load_and_stats(mesh_path)
    T = np.asarray(m.triangles)
    adj = build_face_adjacency(T)
    avg_deg = sum(len(n) for n in adj) / len(adj)
    print(f"Adjacency built. Avg neighbors per face: {avg_deg:.2f}")
    
    m.compute_triangle_normals()
    N = np.asarray(m.triangle_normals)
    clusters = cluster_by_normal(T, N, adj, angle_threshold_degree = 6.0)

    print(f"Cluster found: {len(clusters)}")


    flat_faces = select_planar_clusters(
        m, clusters,
        med_edge=stats["med_edge"],
        max_dev_factor=0.2,     # tighten/loosen if needed
        area_min_factor=50.0,
        rms_factor = 0.20
    )
    
    print(f"Selected planar faces: {len(flat_faces)}")

    m.compute_vertex_normals()
    if not m.has_vertex_colors():
        m.vertex_colors = o3d.utility.Vector3dVector(
            np.zeros((len(m.vertices), 3), dtype=np.float64)
        )
        
    Vcols = np.asarray(m.vertex_colors)
    bone_col = np.array([0.96, 0.94, 0.8])  # light bone
    red_col  = np.array([1.0, 0.0, 0.0])

    Vcols[:] = bone_col

    if len(flat_faces) > 0:
        verts_flat = np.unique(T[flat_faces].ravel())
        Vcols[verts_flat] = red_col
        
    out_colored_path = os.path.join(
        aligned_dir,
        f"{case_number}_{label}_par-cut-aligned_colored.ply"
    )
    o3d.io.write_triangle_mesh(out_colored_path, m, write_vertex_colors=True)
    print(f"Saved colored {label} mesh -> {out_colored_path}")

process_color_mesh(mesh_path_femur, femur_aligned_path, case_number, "fem")
process_color_mesh(mesh_path_tibia, tibia_aligned_path, case_number, "tib")
mesh_repair = input("Is mesh damaged?:")

if mesh_repair.lower() == "yes":
    mesh_colored_path = os.path.join(femur_aligned_path, f"{case_number}_fem_par-cut-aligned_colored.ply")
    mesh = o3d.io.read_triangle_mesh(mesh_colored_path)
    if mesh.is_empty():
        raise RuntimeError("Failed to read mesh")

    # light, safe cleanup
    mesh.remove_degenerate_triangles()
    mesh.remove_duplicated_triangles()
    mesh.remove_duplicated_vertices()
    mesh.remove_non_manifold_edges()
    mesh.compute_vertex_normals()

    m_s = mesh.filter_smooth_laplacian(   # classic Laplacian, but safer after welding
        number_of_iterations=10, lambda_filter=0.09, filter_scope=o3d.geometry.FilterScope.All
    )
    m_s.compute_vertex_normals()
    mesh_repaired_path = os.path.join(femur_aligned_path, f"{case_number}_fem_par-cut-aligned_colored_edited.ply")
    o3d.io.write_triangle_mesh(mesh_repaired_path, m_s)
    print("Mesh repaired.")
else:
    pass

#### TAKE SCREENSHOT #### - THIS TAKES UP THE MAJORITY OF THE TIME WHEN RUN, SO IF YOU OBTAIN THESE IMAGES ONCE, YOU CAN DOC OUT THIS SECTION, AND STILL RUN THE CODE.

pts_files_femur_front = [
    os.path.join(femur_aligned_path, f"{case_number}_fem_par-dis-lat-off-pt-aligned.pts"),
    os.path.join(femur_aligned_path,f"{case_number}_fem_par-dis-med-off-pt-aligned.pts")
]
pts_files_femur_top = [
    os.path.join(femur_aligned_path,f"{case_number}_fem_par-pos-lat-off-pt-aligned.pts"),
    os.path.join(femur_aligned_path,f"{case_number}_fem_par-pos-med-off-pt-aligned.pts")
]

pts_files_tibia = [
    os.path.join(tibia_aligned_path,f"{case_number}_tib_par-lat-sul-pt-aligned.pts"),
    os.path.join(tibia_aligned_path,f"{case_number}_tib_par-med-sul-pt-aligned.pts")
]

tibia_mesh_cut = os.path.join(tibia_aligned_path,f"{case_number}_tib_par-cut-aligned_colored.ply")
tibia_implant =  os.path.join(tibia_aligned_path,f"{case_number}_tib_par-{tibial_raw}-implant-aligned.obj")

tibia_mesh_uncut = os.path.join(tibia_aligned_path,f"{case_number}_tib_par-uncut-aligned.obj")

if mesh_repair.lower() == "yes":    
    femur_mesh_cut = os.path.join(femur_aligned_path,f"{case_number}_fem_par-cut-aligned_colored_edited.ply")
else:
    femur_mesh_cut = os.path.join(femur_aligned_path,f"{case_number}_fem_par-cut-aligned_colored.ply")

femur_implants = [
    os.path.join(femur_aligned_path,f"{case_number}_fem_par-{femoral_raw}-CR-implant-aligned.obj"),
    os.path.join(femur_aligned_path,f"{case_number}_fem_par-{femoral_raw}-PS-implant-aligned.obj"),
]
femur_mesh_uncut = os.path.join(femur_aligned_path,f"{case_number}_fem_par-uncut-aligned.obj")


def set_scene(size=(1000,750)):
    ren = vtk.vtkRenderer()
    ren.GradientBackgroundOff()
    ren.UseDepthPeelingOn(); ren.SetMaximumNumberOfPeels(100); ren.SetOcclusionRatio(0.1) #allows transparent objects to be captured properly, sets the number of layers of transparency, sets how much we're willing to skip small details to save time
    ren.SetBackground(1,1,1) #sets a white bg
    ren.SetBackgroundAlpha(1.0) 
    
    window = vtk.vtkRenderWindow()
    window.OffScreenRenderingOn()
    window.AddRenderer(ren)
    window.SetSize(*size)
    window.SetMultiSamples(0) #prevents smoothing jagged edges so it doesn tlook weird.
    window.SetAlphaBitPlanes(1)#sets up the alpha channel in the render window
    return ren, window

def add_mesh(ren, mesh_paths, color = (0,0,0), opacity=1.0, transparent_bg=False):

    if not isinstance(mesh_paths, list):  # allow single string too
        mesh_paths = [mesh_paths]

    actors = []
    for mesh_path in mesh_paths:
        ext = os.path.splitext(mesh_path)[1].lower()
        if ext == ".ply":
            reader = vtk.vtkPLYReader()
        else:
            reader = vtk.vtkOBJReader()

        reader.SetFileName(mesh_path)
        reader.Update()

        poly = reader.GetOutput()
        n_pts = poly.GetNumberOfPoints()
        b = poly.GetBounds()
        print(f"[LOAD] {os.path.basename(mesh_path)}  points={n_pts}  bounds={b}")
        if n_pts == 0:
            print("  ⚠️  Empty mesh (wrong path or unreadable file)")

        mapper = vtk.vtkPolyDataMapper()
        mapper.SetInputConnection(reader.GetOutputPort())
        pd = poly.GetPointData()
        scalars = pd.GetScalars() if pd else None

        if scalars:
            # Use existing colours
            mapper.ScalarVisibilityOn()
            mapper.SetColorModeToDirectScalars()

            # Replace black [0,0,0] with bone colour
            import vtkmodules.util.numpy_support as nps
            cols = nps.vtk_to_numpy(scalars)
            mask = (cols[:,0]==0) & (cols[:,1]==0) & (cols[:,2]==0)
            cols[mask,:3] = (int(0.89*255), int(0.855*255), int(0.79*255))
            pd.SetScalars(nps.numpy_to_vtk(cols, deep=True, array_type=scalars.GetDataType()))
        else:
            # Fall back to solid bone/implant colour
            mapper.ScalarVisibilityOff()

        actor = vtk.vtkActor()
        actor.SetMapper(mapper)
        p = actor.GetProperty()
        p.SetColor(*color)
        p.SetOpacity(opacity)
        p.SetInterpolationToPhong()
        p.SetAmbient(0.1)
        p.SetDiffuse(0.8)
        p.SetSpecular(0.2)
        p.SetSpecularPower(20)

        ren.AddActor(actor)
        actors.append(actor)

    return actors

def add_points(ren, pts_paths, color = (1,0,0), size = 8.0):
    
    if isinstance(pts_paths, (str, bytes)):
        pts_paths = [pts_paths]
    
    pts = vtk.vtkPoints()
    
    for path in pts_paths:
        with open(path, "r") as f: #reads the pts file
            for line in f: #loops over the lines of the file
                parts = line.strip().split()
                if len(parts) >= 3:
                    try:
                        x, y, z = map(float, parts[:3])
                        if all(map(math.isfinite, (x, y, z))): 
                            pts.InsertNextPoint(x, y, z) #creates a points list for VTK
                    except ValueError:
                            # skip non-numeric/header lines
                                continue
                            
    poly = vtk.vtkPolyData() #creates an empty polydata object, its a data container for geometry
    poly.SetPoints(pts) #adds points into the polydata
    
    vtx = vtk.vtkVertexGlyphFilter() #turns the points into drawable vertices
    vtx.SetInputData(poly)
    vtx.Update()
    
    mapper = vtk.vtkPolyDataMapper() #Creates a polydata mapper, turns a vtkPolyData mesh into GPU draw calls(instructions your computer sends to the graphics card to put something on the screen)
    mapper.SetInputConnection(vtx.GetOutputPort())
    mapper.ScalarVisibilityOff() #ignores colour in the data file

    actor = vtk.vtkActor()
    actor.SetMapper(mapper)
    p = actor.GetProperty() #gets properties of the mesh
    p.SetColor(*color)
    p.SetRepresentationToPoints() #ensures you see the points
    p.SetPointSize(size)

    ren.AddActor(actor)
    
    return actor


def save_img(ren, window, output_png):
    
    ren.ResetCamera() #moves the camera to fit everything into view
    cam = ren.GetActiveCamera() 
    cam.SetParallelProjection(True)
    cam.Zoom(0.9)
    ren.ResetCameraClippingRange() #ensures nothing gets clipped
    
    window.Render()
    window2img = vtk.vtkWindowToImageFilter() #turns the image of Render Window and turns it into a vtkimagedata obj.
    window2img.SetInput(window) #sets the render window to capture
    window2img.SetInputBufferTypeToRGBA() #includes transparency in the image capture
    window2img.ReadFrontBufferOff() #front buffer doesn't have the rendered image
    window2img.Update()
    writer = vtk.vtkPNGWriter() #writes the image to a png file
    writer.SetFileName(output_png)
    writer.SetInputConnection(window2img.GetOutputPort())
    writer.Write()
    print(f"Saved:{os.path.abspath(output_png)}")
    




def set_view_add_mesh(ren, overlay, view_name, type):
    
    ren.RemoveAllViewProps()
    overlay.RemoveAllViewProps()
    overlay.Modified()
    
    
    
    IF YOU WANT TO EDIT THE COLOURS, CHANGE THE VALUES BELOW. 
    
    
    
    
    if type == "top_femur_cut":
        add_mesh(ren, femur_mesh_cut, color = (0.96,0.94,0.8), opacity=1.0, transparent_bg=False)
        add_mesh(ren, femur_implants, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
    elif type == "top_femur_uncut":
        add_mesh(ren, femur_mesh_uncut, color = (0.96,0.94,0.8), opacity=0.5, transparent_bg=False)
        add_mesh(ren, femur_implants, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
        add_points(overlay, pts_files_femur_top, color = (1,0,0), size = 20.0)
    elif type == "top_tibia_cut":
    # Tibia (grab the single actor)
        add_mesh(ren, tibia_mesh_cut, color=(0.96,0.94,0.8), opacity=1.0)
        add_mesh(ren, tibia_implant, color=(0.5,0.5,0.5), opacity=1.0)
    elif type == "top_tibia_uncut":
        add_mesh(ren, tibia_mesh_uncut, color = (0.96,0.94,0.8), opacity=0.5, transparent_bg=False)
        add_mesh(ren, tibia_implant, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
        add_points(overlay, pts_files_tibia, color = (1,0,0), size = 20.0)
    elif type == "front_cut":
        add_mesh(ren, femur_mesh_cut, color = (0.96,0.94,0.8), opacity=1.0, transparent_bg=False)
        add_mesh(ren, femur_implants, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
        add_mesh(ren, tibia_mesh_cut, color = (0.96,0.94,0.8), opacity=1.0, transparent_bg=False)        
        add_mesh(ren, tibia_implant, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
    elif type == "front_uncut":
        add_mesh(ren, femur_mesh_uncut, color = (0.96,0.94,0.8), opacity=0.5, transparent_bg=False)
        add_mesh(ren, femur_implants, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
        add_mesh(ren, tibia_mesh_uncut, color = (0.96,0.94,0.8), opacity=0.5, transparent_bg=False)        
        add_mesh(ren, tibia_implant, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
        add_points(overlay, pts_files_tibia, color = (1,0,0), size = 20.0)
        add_points(overlay, pts_files_femur_front, color = (1,0,0), size = 20.0)
    elif type == "side_cut":
        add_mesh(ren, femur_mesh_cut, color = (0.96,0.94,0.8), opacity=1.0, transparent_bg=False)
        add_mesh(ren, femur_implants, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
        add_mesh(ren, tibia_mesh_cut, color = (0.96,0.94,0.8), opacity=1.0, transparent_bg=False)        
        add_mesh(ren, tibia_implant, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
    elif type == "side_uncut":
        add_mesh(ren, femur_mesh_uncut, color = (0.96,0.94,0.8), opacity=0.5, transparent_bg=False)
        add_mesh(ren, femur_implants, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)
        add_mesh(ren, tibia_mesh_uncut, color = (0.96,0.94,0.8), opacity=0.5, transparent_bg=False)        
        add_mesh(ren, tibia_implant, color = (0.5,0.5,0.5), opacity=1.0, transparent_bg=False)

    overlay.SetActiveCamera(ren.GetActiveCamera())
    ren.ResetCamera()
    cam = ren.GetActiveCamera()
    focal = cam.GetFocalPoint()
    distance = cam.GetDistance()
    
    if view_name == "front":
        cam.SetPosition(focal[0], focal[1], focal[2] + distance) #x-axis remains as is, moves a distance away from the focal point along the y axis

    elif view_name == "left":
        cam.SetPosition(focal[0] - distance, focal[1], focal[2]) #moves along the x_axis

    elif view_name == "right":
        cam.SetPosition(focal[0]+ distance, focal[1], focal[2])
        
    elif view_name == "top_tibia":
        cam.SetPosition(focal[0], focal[1] + distance, focal[2])
        
    elif view_name == "top_femur":
        cam.SetPosition(focal[0], focal[1] - distance, focal[2])

    cam.SetViewUp(0, 1, 0)  # Make Z axis point up
    cam.ParallelProjectionOn()
    ren.ResetCameraClippingRange()
    
ren, window =  set_scene(size=(10000,7500))

overlay = vtk.vtkRenderer()
overlay.SetLayer(1)                    # draw after main
overlay.SetBackgroundAlpha(0.0)        # no background
overlay.SetActiveCamera(ren.GetActiveCamera())  # keep cameras in sync
window.SetNumberOfLayers(2)
window.AddRenderer(overlay)


set_view_add_mesh(ren, overlay, "front", "front_cut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_cut_front_view.png"))

set_view_add_mesh(ren, overlay, "left", "side_cut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_cut_left_view.png"))

set_view_add_mesh(ren, overlay, "right", "side_cut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_cut_right_view.png"))

set_view_add_mesh(ren, overlay, "top_tibia", "top_tibia_cut")
img = save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_cut_top_view_tibia.png"))
img = cv2.imread(os.path.join(case_path,f"{case_number}_cut_top_view_tibia.png"))
img = cv2.rotate(img, cv2.ROTATE_180)
cv2.imwrite(os.path.join(case_path, f"{case_number}_cut_top_view_tibia.png"), img)


set_view_add_mesh(ren, overlay, "top_femur", "top_femur_cut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_cut_top_view_femur.png"))



set_view_add_mesh(ren, overlay, "front", "front_uncut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_uncut_front_view.png"))

set_view_add_mesh(ren, overlay, "left", "side_uncut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_uncut_left_view.png"))

set_view_add_mesh(ren, overlay, "right", "side_uncut")
save_img(ren, window, output_png= os.path.join(case_path,f"{case_number}_uncut_right_view.png"))

set_view_add_mesh(ren, overlay, "top_tibia", "top_tibia_uncut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_uncut_top_view_tibia.png"))
img = cv2.imread(os.path.join(case_path,f"{case_number}_uncut_top_view_tibia.png"))
img = cv2.rotate(img, cv2.ROTATE_180)
cv2.imwrite(os.path.join(case_path, f"{case_number}_uncut_top_view_tibia.png"), img)


set_view_add_mesh(ren, overlay, "top_femur", "top_femur_uncut")
save_img(ren, window, output_png=os.path.join(case_path,f"{case_number}_uncut_top_view_femur.png"))




#### INSERT EXCEL and WORD VALUES ####



slide_number = 2

def delete_slides(ppt, index):
    xml_slides = ppt.slides._sldIdLst
    slides = list(xml_slides)
    ppt.slides._sldIdLst.remove(slides[index])

def updated_mappings(ppt, mappings):
    
    for m in mappings:
        val = str(m["df"].loc[m["Row_Header"], m["Column_Header"]])
        val = f"{float(val):.1f}"
        
        top_min = Cm(m["top_min"])
        top_max = Cm(m["top_max"])
        left_min = Cm(m["left_min"])
        left_max = Cm(m["left_max"])
        
        slide_index = m["slide"] - 1
        should_replace = m.get("replace_if", "number")
        slide = ppt.slides[slide_index]
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                if top_min < shape.top < top_max and left_min < shape.left < left_max:
                    current_text = shape.text.strip()
                    try:
                        float(current_text)  # ✅ Try converting to number
                        shape.text_frame.paragraphs[0].runs[0].text = val
                        print(f"Updated numeric text: '{current_text}' → '{val}'")
                    except ValueError:
                        print(f"⚠️ Skipped non-numeric text: '{current_text}'")

def update_info(case_id, tibia_df, femur_df, word_df, ppt):       
    case_id = str(case_id).strip()
    number_part = case_id[:-1]
    last_number_bit = number_part[:-3]
    letter_part = case_id[-1]
    print(f"📦 Raw case_id: '{case_id}'")
    print(f"🔢 number_part: '{case_id[:-1]}'")
    print(f"🔤 letter_part: '{case_id[-1]}'")
    # Format femoral size from "SZ_30" to "3.0"
    femoral_raw = femur_df.loc["Femoral Size", "Recommended"]
    femoral_size = str(femoral_raw).replace("SZ_", "")
    if femoral_size.isdigit():
        femoral_size = f"{int(femoral_size)/10:.1f}"
    else:
        femoral_size = "NA" 
    tibial_raw = tibia_df.loc["Tibial Size", "Recommended"]
    tibial_size = str(tibial_raw).replace("SZ_", "")
    if tibial_size.isdigit():
        tibial_size = f"{int(tibial_size)/10:.1f}"
    else:
        tibial_size = "NA"       
        
    
    if letter_part == "L":
        letter = "LEFT"
        delete_slides(ppt, 4-1)
        delete_slides(ppt, 5-2)
    elif letter_part == "R":
        letter =  "RIGHT"
        delete_slides(ppt, 2-1)
        delete_slides(ppt, 3-2)
    else:
        letter = "UNKNOWN"
        
    patient_name = word_df.loc[word_df[0] == "Patient Name:", 1].values[0]
    patient_id = word_df.loc[word_df[0] == "Patient ID:", 1].values[0]
    
    replacements = {
        "Case Number #YYMMXXX": (f"Case Number #{str(number_part)}", Pt(20), False, RGBColor(255, 255, 255)),
        "03-SP-YYMMXXX": (f"03-SP-{str(number_part)}", Pt(10), True, RGBColor(0, 0, 0)),
        "AA": (f"{femoral_size}", Pt(12), False, RGBColor(0, 0, 0)),
        "BB": (f"{tibial_size}", Pt(12), False, RGBColor(0, 0, 0)),
        "Knee: D.L": (f"Knee: {letter}", Pt(12), True, RGBColor(0, 0, 0)),
        "Knee: D.R": (f"Knee: {letter}", Pt(12), True, RGBColor(0, 0, 0)),
        "First Initial. Last Name": (f"{patient_name}", Pt(12), False, RGBColor(0, 0, 0)),
        "SN-PN-DDMMYYY": (f"{patient_id}", Pt(12), False, RGBColor(0, 0, 0))
    }
    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    key = run.text.strip()
                    if key in replacements:
                        val, size, bold, color = replacements[key]
                        run.text = val
                        run.font.size = size
                        run.font.bold = bold
                        run.font.color.rgb = color


mappings = [
    {
        "df" : femur_df,
        "Row_Header" : "Femoral Flexion [Mechanical]" ,
        "Column_Header" : "Recommended" ,
        "top_min" :  9.0,
        "top_max" :  11.0,
        "left_min" : 11.0,
        "left_max" : 12.0,
        "slide" : 2
    } ,
    {
        "df" : femur_df,
        "Row_Header" : "Coronal (V/V) Alignment [Anatomical]" ,
        "Column_Header" : "Recommended" ,
        "top_min" :  3.0,
        "top_max" :  4.0,
        "left_min" : 4.0,
        "left_max" : 5.0,
        "slide" : 2
    } ,
    
    {
        "df" : femur_df,
        "Row_Header" : "Rotation to Posterior Condyles" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 12.0,
        "top_max" : 13.0,
        "left_min" : 23.0,
        "left_max" : 24.0,
        "slide" : 2
    } ,
    
    {
        "df" : femur_df,
        "Row_Header" : "Distal Medial Resection" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 10.0,
        "top_max" : 11.0,
        "left_min" : 1.0,
        "left_max" : 2.0,
        "slide" : 2
    } ,
    
    {
        "df" : femur_df,
        "Row_Header" : "Distal Lateral Resection" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 10.0,
        "top_max" : 11.0,
        "left_min" : 7.0,
        "left_max" : 8.0,
        "slide" : 2
    } ,
    
    {  
        "df" : femur_df,
        "Row_Header" : "Posterior Medial Resection", 
        "Column_Header" : "Recommended", 
        "top_min" : 12.0,
        "top_max" : 13.0,
        "left_min" : 20.0,
        "left_max" : 21.0,
        "slide" : 2
    } ,

    {   
        "df" : femur_df,
        "Row_Header" : "Posterior Lateral Resection" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 12.0,
        "top_max" : 13.0,
        "left_min" : 26.0,
        "left_max" : 27.0,
        "slide" : 2
    },

    {
        "df" : tibia_df,
        "Row_Header" : "Posterior Slope [Mechanical]" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 13.0,
        "top_max" : 14.0,
        "left_min" : 12.0,
        "left_max" : 13.0,
        "slide" : 2 
    }, 
    
    {
        "df" : tibia_df,
        "Row_Header" : "Coronal (V/V) Alignment  [Mechanical]" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 13.0,
        "top_max" : 14.0,
        "left_min" : 2.0,
        "left_max" : 3.0,
        "slide" : 2 
    }, 
        
    {
        "df" : tibia_df,
        "Row_Header" : "Medial Resection" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 11.0,
        "top_max" : 12.0,
        "left_min" : 0.0,
        "left_max" : 1.0,
        "slide" : 2
    }, 
            
    {
        "df" : tibia_df,
        "Row_Header" : "Lateral Resection" ,
        "Column_Header" : "Recommended" ,
        "top_min" : 11.0,
        "top_max" : 12.0,
        "left_min" : 8.0,
        "left_max" : 9.0,
        "slide" : 2
    }
    
    ]

#### INSERT X-RAYS ####

pytesseract.pytesseract.tesseract_cmd = "/usr/local/bin/tesseract"

algo_path = os.path.join(base_path, "algo_template.png")
algo_template = cv2.imread(algo_path, cv2.IMREAD_COLOR)


def cm_to_px(cm, dpi=96):
    return int((cm/2.54 * dpi))

def mask_with_template_simple(image, template, scale = 0.4):
    
    # Convert both to grayscale
    img_gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    tpl_gray = cv2.cvtColor(template, cv2.COLOR_BGR2GRAY)
    if scale != 1.0:
        tpl_gray = cv2.resize(
            tpl_gray,
            (int(tpl_gray.shape[1] * scale), int(tpl_gray.shape[0] * scale)),
            interpolation=cv2.INTER_AREA
        )
    # Match template
    result = cv2.matchTemplate(img_gray, tpl_gray, cv2.TM_CCOEFF_NORMED)

    # Find best match
    _, max_val, _, max_loc = cv2.minMaxLoc(result)

    if max_val > 0.5:  # TWEAK IF NEEDED, REDUCE IF NOT IDENTIFIED
        h, w = tpl_gray.shape
        x, y = max_loc
        # Black out that region
        image[y:y+h, x:x+w] = 0

    return image

def fit_image_to_box(img, target_width_px, target_height_px, save_path, img_path=None, top_trim_px = 2):

    if top_trim_px > 0:
        img = img[top_trim_px:, :]
    # Get the original dimensions of the image
    h, w = img.shape[:2]  # h = height (rows), w = width (columns)
    top_bias_ratio = 0
    # Compute a uniform scale factor to resize the image
    # so that it fits within the target dimensions *without distortion*
    # We choose the smaller scaling ratio so the image doesn't exceed the box
    scale = max(target_height_px / h, target_width_px / w)

    # Apply the scale to both dimensions to maintain aspect ratio
    new_w = int(w * scale)  # scaled width
    new_h = int(h * scale)  # scaled height

    # Resize the original image to the new dimensions
    # NOTE: cv2.resize expects (width, height) as (cols, rows)
    resized = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_AREA)
        # Now crop the center to fit exactly
    global algo_template
    if algo_template is not None:
        resized = mask_with_template_simple(resized, algo_template)
    
    label_height = 24
    label_width = 70
    
    side_bias_ratio = 0.7
    total_x_crop = max((new_w - target_width_px), 0)  #how much to crop off from the left edge
    x_crop = int(total_x_crop * side_bias_ratio)
    total_y_crop = new_h - target_height_px 
    crop_top = int(total_y_crop * top_bias_ratio) #how much to crop off the top edge
    y_crop = crop_top
    y_crop = min(y_crop, new_h-target_height_px)
    cropped = resized[y_crop:y_crop + target_height_px, x_crop:x_crop + target_width_px]#sets the size of the new image
    
    cv2.imwrite(save_path, cropped)
    # Return the final, centered, padded image ready for saving or insertion
    return save_path

def insert_xrays(ppt, image_map, slide_index = 0):
    
    slide = ppt.slides[slide_index]
    for i, shape in enumerate(slide.shapes):
        print(f"Index: {i}, Type: {shape.shape_type}, Text: {getattr(shape, 'text', '')}")
    
    for label, info in image_map.items():
        img = cv2.imread(info["img_path"])
        width_px = cm_to_px(info["width"])
        height_px = cm_to_px(info["height"])
        
        processed_path = fit_image_to_box(img, width_px, height_px, info["output_img"], img_path=info["img_path"], top_trim_px= info.get("trim_top", 0))
        
        pic = slide.shapes.add_picture(
                processed_path,
                Cm(info["left"]),
                Cm(info["top"]),
                width = Cm(info["width"]),
                height= Cm(info["height"])
            )
        target_index = 14
        target_shape = slide.shapes[target_index]
        
        spTree = slide.shapes._spTree
        spTree.remove(pic._element)
        target_element = target_shape._element
        target_position = list(spTree).index(target_element)
        
        spTree.insert(target_position, pic._element)

image_map = {
    "femur_front": {
        "img_path": os.path.join(xray_image_path,f"{case_number}_femur_front.png"),
        "output_img": os.path.join(xray_image_path,f"{case_number}_femur_front_processed.jpg"),
        "left": 1.29, "top": 3.63, "width": 7.57, "height": 8.96, "trim_top": 30
    },
    "femur_side": {
        "img_path": os.path.join(xray_image_path,f"{case_number}_femur_side.png"),
        "output_img": os.path.join(xray_image_path,f"{case_number}_femur_side_processed.jpg"),
        "left": 9.17, "top": 3.63, "width": 6.5, "height": 8.96, "trim_top": 45
    },
    "tibia_front": {
        "img_path":  os.path.join(xray_image_path,f"{case_number}_tibia_front.png"),
        "output_img": os.path.join(xray_image_path,f"{case_number}_tibia_front_processed.jpg"),
        "left": 1.29, "top": 12.59, "width": 7.8, "height": 7.65 , "trim_top": 30
    },
    "tibia_side": {
        "img_path": os.path.join(xray_image_path,f"{case_number}_tibia_side.png"),
        "output_img": os.path.join(xray_image_path,f"{case_number}_tibia_side_processed.jpg"),
        "left": 9.26, "top": 12.38, "width": 6.3, "height": 7.8
    }
}

#### INSERT MESH TO PPTX ####



def trim_whitespace(img, white_tol=180, pad_px=0):

    white_tol = max(0, min(int(white_tol), 255)) #it can't be outside 0-255
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY) #converts to grayscale, easier to tell brightness

    # content = darker than threshold (i.e., not background)
    mask = gray < white_tol #where it is less bright than white
    if not np.any(mask):
        return img  # nothing to trim

    ys, xs = np.where(mask) #finds the rows and columns where there is content(darker than white)
    y0, y1 = ys.min(), ys.max() + 1 #adds 1 to the top and bottom row/column because python slicing is exclusive
    x0, x1 = xs.min(), xs.max() + 1
    trimmed = img[y0:y1, x0:x1] #trimmed image is just the bits shown

    if pad_px <= 0:
        return trimmed

    h, w = trimmed.shape[:2]
    top = bottom = left = right = int(pad_px) #adds padding equally

    # pad with white so background stays clean
    return cv2.copyMakeBorder(trimmed, top, bottom, left, right,borderType=cv2.BORDER_CONSTANT, value=(255, 255, 255)) #adds a border


def delete_slides(ppt, index):
    xml_slides = ppt.slides._sldIdLst
    slides = list(xml_slides)
    ppt.slides._sldIdLst.remove(slides[index])
    
def cm_to_px(cm, dpi=96):
    return int((cm/2.54 * dpi))


def crop_to_aspect_cover(img, target_w_cm, target_h_cm, save_path, top_bias_ratio=0.0):
    """
    Crop the image to match the target aspect ratio (target_w_cm : target_h_cm)
    WITHOUT resizing. This is a 'cover' crop: it fills the frame with no padding.
    top_bias_ratio in [0,1]: 0 = center, 1 = keep as much TOP as possible.
    """
    h, w = img.shape[:2] #gets the height and width
    assert h > 0 and w > 0, "Invalid image dims"

    ar_target = float(target_w_cm) / float(target_h_cm)
    ar_img = w / h #gets aspect ratio 

    if ar_img > ar_target:
        # Image is wider than frame → crop width
        crop_w = int(round(h * ar_target))
        x0 = (w - crop_w) // 2
        y0 = 0
        x1 = x0 + crop_w
        y1 = h #crops width wise
    else:
        # Image is taller than frame → crop height (apply top bias)
        crop_h = int(round(w / ar_target))
        extra_h = h - crop_h
        # top_bias_ratio = 0 → center; 1 → align to top
        y0 = int(round(extra_h * (1 - top_bias_ratio) / 2.0))
        y0 = max(0, min(y0, extra_h))  # clamp
        x0 = 0
        x1 = w
        y1 = y0 + crop_h #crops height wise

    cropped = img[y0:y1, x0:x1]
    ok = cv2.imwrite(save_path, cropped)
    assert ok, f"Failed to write {save_path}"
    return save_path

def prepare_for_frame(img, w_cm, h_cm, save_path, top_bias_ratio=0.0, white_tol=245, pad_px=0):
    trimmed = trim_whitespace(img, white_tol=white_tol, pad_px=pad_px)
    return crop_to_aspect_cover(trimmed, w_cm, h_cm, save_path, top_bias_ratio=top_bias_ratio)


def remove_shape(slide, remove_ids):
    # sort in reverse so later deletions don't mess up earlier indices
    for shp in list(slide.shapes):
        if shp.shape_id in remove_ids:
            slide.shapes._spTree.remove(shp.element)

def insert_3D(image_map, side_cut, ppt):
    
    slide_left_uncut = ppt.slides[1]
    slide_left_cut = ppt.slides[2]
    slide_right_uncut = ppt.slides[3]
    slide_right_cut = ppt.slides[4]
    
    if side_cut == "left_uncut":
        slide = slide_left_uncut
        remove_shape(slide, [16, 15, 14, 22])    
    elif side_cut == "left_cut":
        slide = slide_left_cut
        remove_shape(slide, [39, 37, 16, 38])    
    elif side_cut == "right_uncut":
        slide = slide_right_uncut
        remove_shape(slide, [20, 14, 93, 29])  
    elif side_cut == "right_cut":
        slide = slide_right_cut
        remove_shape(slide, [115, 60, 66, 65])
        
    slide_idx = ppt.slides.index(slide)


        
    for i, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"Slide index: {slide_idx}, Index: {i}, Type: {shape.shape_type}, ID: {shape.shape_id})")
        
    for label, info in image_map.items():
        img = cv2.imread(info["img_path"])
        width_px = cm_to_px(info["width"])
        height_px = cm_to_px(info["height"])
        
        processed_path = prepare_for_frame(
    img,
    info["width"], info["height"],
    info["output_img"],
    top_bias_ratio=-0.14,
    pad_px=65,
    white_tol= 250         
)
        
        pic = slide.shapes.add_picture(
                processed_path,
                Cm(info["left"]),
                Cm(info["top"]),
                width = Cm(info["width"]),
                height= Cm(info["height"])
            )
        target_index = 1 #defines the index of the image we want our images to be just behind
        target_shape = slide.shapes[target_index] #gets the shape at that index
        
        spTree = slide.shapes._spTree #places our images at that index
        spTree.remove(pic._element) #removes the image we want our stuff to be behind
        target_element = target_shape._element #gets this thing called XML which helps us reinsert stuff
        target_position = list(spTree).index(target_element) #converts the shape tree into a list of shapes
        
        spTree.insert(target_position, pic._element) #places our image at the right position, essentially making sure its not right at the front
        
    return ppt

image_map_right_uncut = {
    "tibia_top": {
        "img_path": os.path.join(case_path, f"{case_number}_uncut_top_view_tibia.png"),
        "output_img": "save.png",
        "left": 6.86, "top": 14.39, "width": 3.2, "height": 2.29
    },
    "femur_tibia_front": {
        "img_path": os.path.join(case_path,f"{case_number}_uncut_front_view.png"),
        "output_img": "save.png",
        "left": 2.34, "top": 4.47, "width": 5.62, "height": 13.02
    },
    "femur_tibia_side_right": {
        "img_path": os.path.join(case_path, f"{case_number}_uncut_right_view.png"),
        "output_img": "save.png",
        "left": 12.78, "top": 4.47, "width": 5.31, "height": 13.02
    },
    "femur_top": {
        "img_path": os.path.join(case_path, f"{case_number}_uncut_top_view_femur.png"),
        "output_img": "save.png",
        "left": 21.12, "top": 7.55, "width": 5.89, "height": 4.92
    },
}

image_map_left_uncut = {
    "tibia_top": {
        "img_path": os.path.join(case_path, f"{case_number}_uncut_top_view_tibia.png"),
        "output_img": "save.png",
        "left": 6.86, "top": 14.39, "width": 3.2, "height": 2.29
    },
    "femur_tibia_front": {
        "img_path": os.path.join(case_path, f"{case_number}_uncut_front_view.png"),
        "output_img": "save.png",
        "left": 2.34, "top": 4.47, "width": 5.62, "height": 13.02
    },
    "femur_tibia_side_left": {
        "img_path": os.path.join(case_path,f"{case_number}_uncut_left_view.png"),
        "output_img": "save.png",
        "left": 12.63, "top": 4.94, "height": 12.31, "width": 5.72
},
    "femur_top": {
        "img_path": os.path.join(case_path,f"{case_number}_uncut_top_view_femur.png"),
        "output_img": "save.png",
        "left": 21.12, "top": 7.55, "width": 5.89, "height": 4.92
    }
}

image_map_right_cut = {
    "tibia_top": {
        "img_path": os.path.join(case_path, f"{case_number}_cut_top_view_tibia.png"),
        "output_img": "save.png",
        "left": 6.86, "top": 14.39, "width": 3.2, "height": 2.29
    },
    "femur_tibia_front": {
        "img_path": os.path.join(case_path,f"{case_number}_cut_front_view.png"),
        "output_img": "save.png",
        "left": 2.34, "top": 4.47, "width": 5.62, "height": 13.02
    },
    "femur_tibia_side_left": {
        "img_path": os.path.join(case_path,f"{case_number}_cut_right_view.png"),
        "output_img": "save.png",
        "left": 12.63, "top": 4.94, "height": 12.31, "width": 5.72
},
    "femur_top": {
        "img_path": os.path.join(case_path,f"{case_number}_cut_top_view_femur.png"),
        "output_img": "save.png",
        "left": 21.12, "top": 7.55, "width": 5.89, "height": 4.92
    }
}
image_map_left_cut = {
    "tibia_top": {
        "img_path": os.path.join(case_path, f"{case_number}_cut_top_view_tibia.png"),
        "output_img": "save.png",
        "left": 6.86, "top": 14.39, "width": 3.2, "height": 2.29
    },
    "femur_tibia_front": {
        "img_path": os.path.join(case_path,f"{case_number}_cut_front_view.png"),
        "output_img": "save.png",
        "left": 2.34, "top": 4.47, "width": 5.62, "height": 13.02
    },
    "femur_tibia_side_left": {
        "img_path": os.path.join(case_path,f"{case_number}_cut_left_view.png"),
        "output_img": "save.png",
        "left": 12.63, "top": 4.94, "height": 12.31, "width": 5.72
},
    "femur_top": {
        "img_path": os.path.join(case_path,f"{case_number}_cut_top_view_femur.png"),
        "output_img": "save.png",
        "left": 21.12, "top": 7.55, "width": 5.89, "height": 4.92
    }

}



ppt = Presentation(ppt_path)

insert_3D(
    image_map_right_uncut,
    side_cut = "right_uncut",
    ppt = ppt
)
insert_3D(
    image_map_left_uncut,
    side_cut = "left_uncut",
    ppt= ppt
)

insert_3D(
    image_map_right_cut,
    side_cut = "right_cut",
    ppt= ppt
)

insert_3D(
    image_map_left_cut,
    side_cut = "left_cut",
    ppt= ppt
)
updated_mappings(ppt, mappings)

update_info(case_id, tibia_df, femur_df, word_df, ppt)

insert_xrays(ppt, image_map, slide_index = 0)


ppt.save(os.path.join(case_path,f"{case_number}-surgical-preop-plan.pptx"))